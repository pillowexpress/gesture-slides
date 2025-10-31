import argparse
import sys
import json
import hashlib
from pathlib import Path
from io import BytesIO

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:
    print("ERROR: python-pptx not installed. Run: pip install -r requirements.txt", file=sys.stderr)
    sys.exit(2)

try:
    from PIL import Image
except Exception:
    print("ERROR: Pillow not installed. Run: pip install -r requirements.txt", file=sys.stderr)
    sys.exit(2)

# Namespaces for XML search
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

ASSET_PREFIX = '/slides/images/'  # absolute path used by web app


def _safe_lines_from_shape(shape):
    lines = []
    try:
        if hasattr(shape, 'text') and shape.text:
            for line in shape.text.splitlines():
                s = line.strip()
                if s:
                    lines.append(s)
    except Exception as e:
        print(f"WARN: text extract failed: {e}", file=sys.stderr)
    return lines


def _iter_all_shapes(shapes):
    for shp in shapes:
        try:
            yield shp
            if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub in _iter_all_shapes(shp.shapes):
                    yield sub
        except Exception as e:
            print(f"WARN: shape iterate error: {e}", file=sys.stderr)
            continue


def _save_blob_as_image(blob: bytes, preferred_ext: str, out_path: Path) -> str | None:
    try:
        if preferred_ext.lower() in ('png', 'jpg', 'jpeg', 'webp'):
            with open(out_path, 'wb') as f:
                f.write(blob)
            return out_path.name
        with Image.open(BytesIO(blob)) as im:
            im = im.convert('RGBA') if im.mode in ('P', 'LA') else im.convert('RGB') if im.mode not in ('RGB','RGBA') else im
            out_png = out_path.with_suffix('.png')
            im.save(out_png, format='PNG')
            return out_png.name
    except Exception as e:
        print(f"WARN: image save failed: {e}", file=sys.stderr)
        try:
            with Image.open(BytesIO(blob)) as im:
                out_png = out_path.with_suffix('.png')
                im.save(out_png, format='PNG')
                return out_png.name
        except Exception as e2:
            print(f"WARN: fallback save failed: {e2}", file=sys.stderr)
            return None


def _extract_blip_images_from_element(part, element, out_dir: Path, slide_index: int, start_idx: int, seen_hashes: set[str]) -> tuple[list[str], int]:
    results: list[str] = []
    img_i = start_idx
    for blip in element.findall(f'.//{{{NS_A}}}blip'):
        rId = blip.get(f'{{{NS_R}}}embed')
        if not rId:
            continue
        try:
            rel = part._rels.get(rId)
            if rel is None or rel.target_part is None:
                continue
            blob = rel.target_part.blob
            h = hashlib.sha1(blob).hexdigest()
            if h in seen_hashes:
                continue
            seen_hashes.add(h)
            filename = getattr(rel.target_part, 'filename', '')
            ext = filename.split('.')[-1].lower() if '.' in filename else 'png'
            if ext not in ('png','jpg','jpeg','webp'):
                ext = 'png'
            rel_name = f"{slide_index:03d}_{img_i}.{ext}"
            out_path = out_dir / rel_name
            saved_name = _save_blob_as_image(blob, ext, out_path)
            if saved_name:
                results.append(saved_name)
                img_i += 1
        except Exception as e:
            print(f"WARN: rel resolve failed: {e}", file=sys.stderr)
            continue
    return results, img_i


def _extract_images(slide, out_dir: Path, slide_index: int) -> list[str]:
    results: list[str] = []
    seen_hashes: set[str] = set()
    img_i = 1
    part = slide.part

    for shp in _iter_all_shapes(slide.shapes):
        try:
            names, img_i = _extract_blip_images_from_element(part, shp._element, out_dir, slide_index, img_i, seen_hashes)
            results.extend(names)
        except Exception as e:
            print(f"WARN: element blip scan failed: {e}", file=sys.stderr)
            continue

    try:
        bg = slide._element.find(f'.//{{{NS_A}}}blipFill')
        if bg is not None:
            names, img_i = _extract_blip_images_from_element(part, bg, out_dir, slide_index, img_i, seen_hashes)
            results.extend(names)
    except Exception as e:
        print(f"WARN: bg scan failed: {e}", file=sys.stderr)

    return results


def export_pptx(pptx_path: Path, images_dir: Path, md_path: Path, json_path: Path):
    prs = Presentation(str(pptx_path))
    images_dir.mkdir(parents=True, exist_ok=True)

    slides_manifest = []
    md_lines: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            texts.extend(_safe_lines_from_shape(shape))
        title = texts[0] if texts else f"Slide {i}"
        bullets = texts[1:] if texts else []

        images = _extract_images(slide, images_dir, i)

        slides_manifest.append({
            "index": i,
            "title": title,
            "bullets": bullets,
            "images": [ASSET_PREFIX + name for name in images]
        })

        md_lines.append(f"# {title}")
        for b in bullets:
            md_lines.append(f"- {b}")
        for name in images:
            md_lines.append("")
            md_lines.append(f"![Slide {i}]({ASSET_PREFIX}{name})")
        md_lines.append("")
        md_lines.append("---")

    if md_lines and md_lines[-1] == "---":
        md_lines = md_lines[:-1]

    md_path.parent.mkdir(parents=True, exist_ok=True)
    md_path.write_text("\n".join(md_lines), encoding='utf-8')

    json_path.parent.mkdir(parents=True, exist_ok=True)
    json_path.write_text(json.dumps({"slides": slides_manifest}, ensure_ascii=False, indent=2), encoding='utf-8')


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--in', dest='inp', required=True)
    parser.add_argument('--out', dest='out', required=True)
    parser.add_argument('--md', dest='md', required=True)
    parser.add_argument('--json', dest='json', required=True)
    args = parser.parse_args()

    in_path = Path(args.inp)
    out_dir = Path(args.out)
    md_path = Path(args.md)
    json_path = Path(args.json)

    if not in_path.exists():
        print("ERROR: input file not found", file=sys.stderr)
        sys.exit(1)

    export_pptx(in_path, out_dir, md_path, json_path)
    print("OK")


if __name__ == '__main__':
    main()
